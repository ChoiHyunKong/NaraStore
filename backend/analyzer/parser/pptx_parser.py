"""
PPTX 파일 파싱
python-pptx를 사용한 텍스트 추출
"""
from pptx import Presentation
from typing import Dict, List
from backend.utils.logger import logger
from backend.utils.error_handler import error_handler


class PPTXParser:
    """PPTX 파서 클래스"""
    
    @staticmethod
    def extract_text(file_path: str) -> tuple[bool, str | Dict]:
        """
        PPTX 파일에서 텍스트 추출
        
        Args:
            file_path: PPTX 파일 경로
            
        Returns:
            (성공 여부, 추출된 텍스트 또는 에러 메시지)
        """
        try:
            logger.info(f"PPTX 파싱 시작: {file_path}")
            
            prs = Presentation(file_path)
            
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- 슬라이드 {slide_num} ---\n"
                
                # 슬라이드 내 모든 도형의 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                    
                    # 표 내용 추출
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            slide_text += row_text + "\n"
                
                # 노트 추출
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    notes_text = notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        slide_text += f"\n[노트]\n{notes_text}\n"
                
                slides_text.append(slide_text)
            
            full_text = "\n".join(slides_text)
            
            result = {
                "text": full_text.strip(),
                "total_slides": len(prs.slides)
            }
            
            logger.info(f"PPTX 파싱 완료: {len(prs.slides)}개 슬라이드")
            return True, result
            
        except Exception as e:
            return error_handler.handle_parsing_error(e, "PPTX")
    
    @staticmethod
    def extract_text_by_slide(file_path: str) -> tuple[bool, List[str] | str]:
        """
        슬라이드별로 텍스트 추출
        
        Returns:
            (성공 여부, 슬라이드별 텍스트 리스트 또는 에러 메시지)
        """
        try:
            prs = Presentation(file_path)
            slides_text = []
            
            for slide in prs.slides:
                text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                slides_text.append(text)
            
            return True, slides_text
            
        except Exception as e:
            return error_handler.handle_parsing_error(e, "PPTX")


# 전역 인스턴스
pptx_parser = PPTXParser()
